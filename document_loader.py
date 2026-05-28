"""
DocumentLoader — 5-stage structured pipeline for document ingestion.

Pipeline: parse → tables → vision → text_chunk → metadata

For supported rich formats (PDF, DOCX, PPTX, XLSX, images):
  Stage 1  (parse)     — Docling element walk: headings, paragraphs, tables, pictures
  Stage 2  (tables)    — table classification + LLM narration
  Stage 3  (vision)    — 3-pass Ollama VLM image summarization
  Stage 4  (chunking)  — RCTS chunking with contextual prefix (CJK-aware)
  Stage 5  (metadata)  — language detection, SHA-256 hash, chunk_type assignment

For passthrough formats (audio/video, VTT, HTML, MD, ADOC, TEX):
  Docling → Markdown export, wrapped as a single text chunk.

Returns
-------
load()           → list[dict]              (chunk list)
load_directory() → dict[str, list[dict]]   (filename → chunks)
"""
from __future__ import annotations

import base64
import hashlib
import io
import json
import logging
import os
import re
import tempfile
import warnings
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Extension categories
# ---------------------------------------------------------------------------

# Formats that go through the full 5-stage pipeline
PIPELINE_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx",
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp",
}

# Formats that use Docling→Markdown passthrough (single chunk output)
PASSTHROUGH_EXTENSIONS = {
    ".md", ".markdown", ".adoc", ".asciidoc",
    ".tex", ".html", ".htm", ".xhtml",
    ".csv", ".txt", ".text", ".vtt",
    ".wav", ".mp3", ".m4a", ".aac", ".ogg", ".flac",
    ".mp4", ".avi", ".mov",
}

SUPPORTED_EXTENSIONS = PIPELINE_EXTENSIONS | PASSTHROUGH_EXTENSIONS

# ---------------------------------------------------------------------------
# Docling element label sets
# ---------------------------------------------------------------------------

HEADING_LABELS = {"section_header", "heading", "title"}
PARAGRAPH_LABELS = {"paragraph", "text", "caption", "footnote"}
LIST_LABELS = {"list_item"}
TABLE_LABELS = {"table"}
PICTURE_LABELS = {"picture"}
TEXT_TYPES = {"paragraph", "heading", "list_item"}

# ---------------------------------------------------------------------------
# Table detection / classification helpers (from 02_table.py)
# ---------------------------------------------------------------------------

_BOX_DRAWING_HEAVY_RE = re.compile(r'[│║┌┐└┘├┤┬┴┼═╔╗╚╝╠╣╦╩╬─]{3,}')


def _has_heavy_box_drawing(text: str) -> bool:
    if _BOX_DRAWING_HEAVY_RE.search(text):
        return True
    box_chars = sum(1 for ch in text if ch in '│║┌┐└┘├┤┬┴┼═╔╗╚╝╠╣╦╩╬─')
    return len(text) > 0 and box_chars / len(text) > 0.15


def _estimate_table_dimensions(text: str) -> Tuple[int, int]:
    rows = 0
    max_cols = 0
    for line in text.splitlines():
        stripped = line.strip()
        if re.match(r'^[\|\+\-\+]+$', stripped.replace(' ', '')):
            continue
        if '|' in stripped:
            cells = [c for c in stripped.split('|') if c.strip()]
            if cells:
                rows += 1
                max_cols = max(max_cols, len(cells))
    return rows, max_cols


def _ocr_difficulty(rows: int, cols: int) -> str:
    if cols >= 7 or rows >= 16:
        return "HARD"
    if cols >= 4 or rows >= 9:
        return "MEDIUM"
    return "EASY"


def _classify_table_type(markdown_text: str, rows: int, cols: int) -> str:
    content_lines = [
        l.strip() for l in markdown_text.splitlines()
        if '|' in l.strip()
        and not (re.match(r'^[|\-:+ ]+$', l.strip()) and '-' in l)
    ]
    if not content_lines:
        return "general"
    header = content_lines[0].lower()
    faq_header_kw = ("question", "answer", "q&a", "faq", "| q |", "| a |")
    if any(kw in header for kw in faq_header_kw):
        return "faq"
    data_lines = content_lines[1:]
    if data_lines:
        q_ratio = sum(1 for l in data_lines if "?" in l) / len(data_lines)
        if q_ratio > 0.3:
            return "faq"
    if cols >= 3 and data_lines:
        all_cells: List[str] = []
        for l in data_lines:
            cells = [c.strip() for c in l.split("|") if c.strip()]
            all_cells.extend(cells[1:])
        if all_cells:
            avg_len = sum(len(c) for c in all_cells) / len(all_cells)
            if avg_len < 30:
                return "spec"
    return "general"


def _parse_faq_rows(markdown_text: str) -> Tuple[List[str], List[List[str]]]:
    headers: List[str] = []
    data_rows: List[List[str]] = []
    for line in markdown_text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        is_row_start = stripped.startswith('|')
        if is_row_start and re.match(r'^[|\-:+ ]+$', stripped) and '-' in stripped:
            continue
        if is_row_start:
            cells = [c.strip() for c in stripped.split('|')]
            if cells and cells[0] == '':
                cells = cells[1:]
            if cells and cells[-1] == '':
                cells = cells[:-1]
            if not cells:
                continue
            if not headers:
                headers = cells
            else:
                while len(cells) < len(headers):
                    cells.append('')
                data_rows.append(cells[:len(headers)])
        else:
            if not data_rows:
                continue
            content = stripped.rstrip('| ').strip()
            if not content:
                continue
            row = data_rows[-1]
            for i in range(len(row) - 1, -1, -1):
                if row[i]:
                    row[i] = row[i] + '\n' + content
                    break
    return headers, data_rows


_Q_HEADERS = {"q", "question", "questions", "問題"}
_A_HEADERS = {"a", "answer", "answers", "答案", "回答", "solution", "solutions", "description"}


def _faq_row_to_text(headers: List[str], row: List[str]) -> str:
    q_idx: Optional[int] = None
    a_idx: Optional[int] = None
    for i, h in enumerate(headers):
        hl = h.lower().strip()
        if hl in _Q_HEADERS:
            q_idx = i
        elif hl in _A_HEADERS:
            a_idx = i
    if q_idx is None or a_idx is None:
        if len(headers) == 2:
            q_idx, a_idx = 0, 1
        elif len(headers) >= 3:
            q_idx, a_idx = 1, 2
        else:
            q_idx, a_idx = 0, 0
    q_text = row[q_idx].strip() if q_idx is not None and q_idx < len(row) else ""
    a_text = row[a_idx].strip() if a_idx is not None and a_idx < len(row) else ""
    if q_text and a_text:
        return f"Q: {q_text}\nA: {a_text}"
    if a_text:
        return a_text
    return ""


# ---------------------------------------------------------------------------
# Narration preamble stripping (shared by table + vision)
# ---------------------------------------------------------------------------

_PREAMBLE_PREFIXES = (
    "here's", "here is", "sure", "certainly", "of course",
    "below is", "below are", "the following",
)


def _strip_preamble(text: str) -> str:
    lines = text.splitlines()
    while lines:
        first = lines[0].strip().lower()
        if any(first.startswith(p) for p in _PREAMBLE_PREFIXES):
            lines.pop(0)
        else:
            break
    return "\n".join(lines).strip()


# ---------------------------------------------------------------------------
# Vision type classification (from 03_vision.py)
# ---------------------------------------------------------------------------

_STRUCTURED_PROMPTS: Dict[str, str] = {
    "flowchart": (
        "This is a flowchart. Read every piece of text in the image.\n"
        "Output two sections:\nNODES:\n  - [exact text inside each shape or box]\n"
        "CONNECTIONS:\n  - [source node text] --> [destination node text] (label: [arrow label or none])\n"
        "List every node and every arrow. Do not skip any."
    ),
    "architecture": (
        "This is an architecture or system diagram. Read every piece of text in the image.\n"
        "Output two sections:\nCOMPONENTS:\n  - [component/box label] (group: [surrounding group label if any])\n"
        "CONNECTIONS:\n  - [source label] --> [destination label] (label: [connection label or none])\n"
        "List every component and every connecting line or arrow. Do not skip any."
    ),
    "parts": (
        "This is a parts or assembly diagram. Read every piece of text in the image.\n"
        "For each labelled part or component:\n  - [part number if any] [part name] (location: [where it is])\n"
        "Then list any callout lines:\n  - [label] --> [component it points to]\nDo not skip any label."
    ),
    "chart": (
        "This is a chart or graph. Read every piece of text in the image.\n"
        "Output: TITLE, X-AXIS, Y-AXIS, LEGEND, DATA. List every data point visible."
    ),
    "table": (
        "This is a table. Read every piece of text in the image.\n"
        "HEADERS: [col1] | [col2] | ...\n  - [col1 value] | [col2 value] | ...\n"
        "Preserve every row and every cell. Use '(empty)' for blank cells."
    ),
    "timeline": (
        "This is a timeline. For each milestone:\n  - [date or time]: [event description]\n"
        "List every entry in chronological order."
    ),
    "infographic": (
        "This is an infographic. For each section or panel, output a header line, "
        "then list every text item as:\n  - [text item]\nPreserve numbers and statistics."
    ),
    "schematic": (
        "This is a schematic or circuit diagram.\n"
        "COMPONENTS:\n  - [ref] [type] [value]\nCONNECTIONS:\n  - [net] connects [ref] to [ref]\n"
        "List every component and labelled net."
    ),
    "logo": (
        "This is a logo or brand mark.\nBRAND NAME: [primary name]\n"
        "TAGLINE: [slogan or '(none)']\nOTHER TEXT: [remaining visible text]"
    ),
    "badge": (
        "This is a certification badge.\nCERTIFICATION: [name and number]\n"
        "ISSUING BODY: [body or '(not visible)']\nOTHER TEXT: [remaining text]"
    ),
    "photo": (
        "This is a product or object photograph.\n"
        "SUBJECT: [one sentence]\nVISIBLE TEXT:\n  - [each text string]\nLABELS/STICKERS: [describe labels]"
    ),
}

_DEFAULT_STRUCTURED_PROMPT = (
    "Read every visible text element in this image.\n"
    "Group related items that appear in the same region, box, or section.\n"
    "For each group write a header line, then list items as '  - [text]'.\n"
    "Preserve spatial relationships. Do not skip any text."
)

_COMPLEX_TYPES: frozenset = frozenset({
    "flowchart", "architecture", "chart", "table", "timeline", "infographic", "schematic",
})
_MIN_AREA_FOR_COMPLEX: int = 50_000

_CAPTION_PROMPT = (
    "What type of image is this? Reply in one sentence only. "
    "Choose the single best match: "
    "'This is a system architecture diagram.', 'This is a UI screenshot.', "
    "'This is a flowchart.', 'This is a parts diagram.', "
    "'This is a bar chart.', 'This is a line chart.', 'This is a pie chart.', "
    "'This is a table.', 'This is a timeline.', 'This is an infographic.', "
    "'This is a schematic diagram.', 'This is a logo.', "
    "'This is a certification badge.', 'This is a product photo.', 'This is a map.'"
)

_SYNTHESIS_PROMPT = """\
An image has been analysed. Here are the results:

Image type: {caption}

<structured_text>
{ocr_deduped}
</structured_text>

Rewrite the above as a clean structured description:
1. First line: write only the image type as plain text. No '#' prefix, no markdown.
2. Preserve ALL groups, columns, sections, and connection blocks exactly.
3. Under each group, list every item as a bullet (  * item).
4. For connections, preserve as '  * [source] --> [destination]'.
5. Fix only clear OCR noise but preserve all real text.
6. Do NOT omit any item. Do NOT add information not in <structured_text>.
7. Output ONLY the formatted description. No meta-commentary.\
"""

_OCR_ECHO_PHRASES = (
    "read all text in this image", "output each text item",
    "no explanations", "extract every visible text string",
    "one item per line", "nothing else",
)


def _is_ocr_echo(text: str) -> bool:
    lowered = text.lower()
    return sum(phrase in lowered for phrase in _OCR_ECHO_PHRASES) >= 2


def _dedup_ocr(ocr_text: str) -> str:
    lines = [line for line in ocr_text.splitlines() if line.strip()]
    seen: List[str] = []
    for line in lines:
        if not seen or line != seen[-1]:
            seen.append(line)
    return "\n".join(seen)


def _classify_vision_caption(caption: str) -> str:
    low = caption.lower()
    if any(w in low for w in ("flowchart", "flow chart", "flow diagram", "process diagram", "workflow")):
        return "flowchart"
    if any(w in low for w in ("architecture", "system diagram", "network diagram", "service diagram", "component diagram")):
        return "architecture"
    if any(w in low for w in ("parts", "assembly", "exploded")):
        return "parts"
    if any(w in low for w in ("infographic",)):
        return "infographic"
    if any(w in low for w in ("bar chart", "line chart", "pie chart", "histogram", "scatter plot", "chart", "graph")):
        return "chart"
    if any(w in low for w in ("schematic", "circuit", "wiring diagram", "electrical diagram")):
        return "schematic"
    if any(w in low for w in ("logo", "brand mark", "brand logo")):
        return "logo"
    if any(w in low for w in ("certification badge", "rating badge", "badge", "certification", "seal", "stamp")):
        return "badge"
    if any(w in low for w in ("product photo", "photograph", "photo", "product image")):
        return "photo"
    return "default"


def _size_guard(image_path: str, type_key: str) -> str:
    if type_key not in _COMPLEX_TYPES:
        return type_key
    try:
        from PIL import Image as _PILImage
        with _PILImage.open(image_path) as img:
            w, h = img.size
            area = w * h
        if area < _MIN_AREA_FOR_COMPLEX:
            logger.info("Size guard: %dx%d (%d px²) too small for '%s' → 'default'", w, h, area, type_key)
            return "default"
    except Exception as exc:
        logger.warning("Size guard: could not read image dimensions: %s", exc)
    return type_key


# ---------------------------------------------------------------------------
# Chunking helpers (from 04_text_chunk.py)
# ---------------------------------------------------------------------------

def _build_context_prefix(doc_stem: str, section_title: Optional[str], page_number: Any) -> str:
    parts = [f"Document: {doc_stem}"]
    if section_title:
        parts.append(f"Section: {section_title}")
    if page_number is not None:
        parts.append(f"Page: {page_number}")
    return "[" + " | ".join(parts) + "]"


def _build_rcts(chunk_size: int, chunk_overlap: int):
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        return RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=[
                "\n\n", "\n",
                "\u3002",  # 。
                ".", "\uff01", "!", "\uff1f", "?",
                "\uff1b", ";", "\uff1a", ":", " ", "",
            ],
        )
    except ImportError:
        logger.warning("langchain_text_splitters not available; using simple character split")
        return None


def _preliminary_type(element: Dict) -> str:
    if element.get("type") == "picture":
        return "picture"
    if element.get("type") == "list_item" and element.get("ordered"):
        return "procedure_step"
    return "text"


def _group_elements(elements: List[Dict], vision_map: Optional[Dict[str, str]] = None) -> List[List[Dict]]:
    vision_map = vision_map or {}
    groups: List[List[Dict]] = []
    current: List[Dict] = []

    def flush():
        if current:
            groups.append(list(current))
            current.clear()

    for el in elements:
        t = el.get("type")
        if t == "picture":
            flush()
            pic_text = vision_map.get(el["id"], "").strip()
            if pic_text:
                groups.append([{**el, "text": pic_text}])
            continue
        if t not in TEXT_TYPES:
            flush()
            continue
        if t == "heading":
            flush()
            groups.append([el])
            continue
        same_section = (
            current and
            current[-1].get("section_path") == el.get("section_path") and
            current[-1].get("page") == el.get("page")
        )
        if not same_section:
            flush()
        current.append(el)

    flush()
    return groups


def _merge_short_text_chunks(chunks: List[Dict]) -> List[Dict]:
    NON_MERGE_TYPES = {"procedure_step"}
    doc_stem = chunks[0]["chunk_id"].rsplit("_", 1)[0] if chunks else "doc"

    def _is_mergeable(c: Dict) -> bool:
        if c.get("preliminary_type") in NON_MERGE_TYPES:
            return False
        t = c.get("text", "").strip()
        return len(t) < 80 and "\n" not in t

    merged: List[Dict] = []
    i = 0
    merge_counter = 0
    while i < len(chunks):
        c = chunks[i]
        if not _is_mergeable(c):
            merged.append(c)
            i += 1
            continue
        group = [c]
        j = i + 1
        while j < len(chunks):
            nxt = chunks[j]
            if not _is_mergeable(nxt):
                break
            if nxt.get("page_number") != c.get("page_number"):
                break
            if nxt.get("section_title") != c.get("section_title"):
                break
            combined = (
                sum(len(x.get("text", "").strip()) for x in group)
                + len(" | ") * len(group)
                + len(nxt.get("text", "").strip())
            )
            if combined >= 400:
                break
            group.append(nxt)
            j += 1
        if len(group) == 1:
            merged.append(c)
        else:
            combined_text = " | ".join(x.get("text", "").strip() for x in group)
            all_eids = []
            for g in group:
                all_eids.extend(g.get("source_element_ids", []))
            merge_counter += 1
            merged.append({
                "chunk_id": f"{doc_stem}_merged_{merge_counter}",
                "source_element_ids": all_eids,
                "text": combined_text,
                "page_number": c.get("page_number"),
                "section_title": c.get("section_title"),
                "preliminary_type": "text",
            })
        i = j
    return merged


def _merge_warning_headers(chunks: List[Dict]) -> List[Dict]:
    WARNING_RE = re.compile(r"\b(DANGER|WARNING|CAUTION|NOTICE|NOTE)\b")
    result: List[Dict] = []
    i = 0
    while i < len(chunks):
        chunk = chunks[i]
        text = chunk.get("text", "").strip()
        is_short_warning = len(text) < 50 and WARNING_RE.search(text) is not None
        if is_short_warning and i + 1 < len(chunks):
            nxt = chunks[i + 1]
            if chunk.get("page_number") == nxt.get("page_number"):
                merged = dict(nxt)
                merged["text"] = text + "\n" + nxt["text"].strip()
                merged["preliminary_type"] = nxt.get("preliminary_type", "text")
                merged["source_element_ids"] = (
                    chunk.get("source_element_ids", []) + nxt.get("source_element_ids", [])
                )
                result.append(merged)
                i += 2
                continue
        result.append(chunk)
        i += 1
    return result


# ---------------------------------------------------------------------------
# Metadata helpers (from 05_metadata.py)
# ---------------------------------------------------------------------------

_LANG_REMAP = {"zh-cn": "ZH-TW", "zh-tw": "ZH-TW", "zh": "ZH-TW"}


def _detect_language(text: str) -> str:
    try:
        from langdetect import detect
        code = detect(text)
        lower = code.lower()
        return _LANG_REMAP.get(lower, code.upper())
    except Exception:
        return "EN"


def _chunk_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


# ---------------------------------------------------------------------------
# Narration table prompt
# ---------------------------------------------------------------------------

_NARRATE_PROMPT_TMPL = """\
Convert the markdown table below into concise natural language sentences.
Rules:
- Output ONLY the sentences. No preamble, no explanation, no heading.
- Preserve every value exactly as written.
- Do not add information not present in the table.
Section context: {section_title}

{markdown_text}"""


# ---------------------------------------------------------------------------
# DocumentLoader
# ---------------------------------------------------------------------------


class DocumentLoader:
    """
    Loads documents using a 5-stage structured pipeline.

    For rich formats (PDF, DOCX, PPTX, XLSX, images), the pipeline:
      Stage 1 — Docling element walk (headings, paragraphs, tables, pictures)
      Stage 2 — Table classification + Ollama LLM narration
      Stage 3 — 3-pass Ollama VLM image summarization
      Stage 4 — RCTS text chunking with contextual prefix
      Stage 5 — Language detection, SHA-256 hash, chunk_type assignment

    For passthrough formats (audio, video, VTT, HTML, MD, ADOC, TEX):
      Docling → Markdown, wrapped as a single text chunk.

    Parameters
    ----------
    ollama_base_url:
        Base URL of the Ollama server (used for both VLM and text LLM).
    vlm_model:
        Ollama VLM model name for image summarization.
    request_timeout:
        HTTP timeout in seconds for Ollama calls.
    """

    def __init__(
        self,
        ollama_base_url: str = "http://10.168.3.58:8088",
        vlm_model: str = "llava:7b",
        request_timeout: float = 60.0,
    ) -> None:
        self._base_url = ollama_base_url.rstrip("/")
        self._vlm_model = vlm_model
        self._timeout = request_timeout
        # Text model for table narration: INDEX_MODEL env var, falling back to vlm_model
        self._text_model = os.environ.get("INDEX_MODEL", vlm_model)
        # Chunk parameters from env
        self._chunk_size = int(os.environ.get("CHUNK_SIZE", "800"))
        self._chunk_overlap = int(os.environ.get("CHUNK_OVERLAP", "80"))
        self._converter = self._build_converter()

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def load(self, file_path: str | Path) -> List[Dict]:
        """Convert *file_path* through the pipeline and return chunk list."""
        path = Path(file_path)
        ext = path.suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file extension: {ext}")
        if ext in PASSTHROUGH_EXTENSIONS:
            return self._load_passthrough(path)
        return self._load_pipeline(path)

    def load_directory(
        self,
        directory: str | Path,
        files_filter: Optional[set] = None,
    ) -> Dict[str, List[Dict]]:
        """Walk *directory*, call ``load()`` on each supported file.

        Parameters
        ----------
        directory:
            Root directory to walk recursively.
        files_filter:
            Optional set of basenames. When provided, only matching files are loaded.

        Returns
        -------
        dict[str, list[dict]]
            Mapping of ``filename → chunk list``.
        """
        directory = Path(directory)
        results: Dict[str, List[Dict]] = {}
        for file_path in directory.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                if files_filter is not None and file_path.name not in files_filter:
                    continue
                try:
                    results[file_path.name] = self.load(file_path)
                    logger.info("Loaded %s (%d chunks)", file_path.name, len(results[file_path.name]))
                except Exception as exc:
                    logger.warning("Failed to load %s: %s", file_path.name, exc)
        return results

    # ------------------------------------------------------------------
    # Passthrough path (single-chunk output for simple formats)
    # ------------------------------------------------------------------

    def _load_passthrough(self, path: Path) -> List[Dict]:
        """Use Docling→Markdown export, then split into chunks respecting CHUNK_SIZE."""
        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", message="Palette images with Transparency",
                                    category=UserWarning, module="PIL")
            result = self._converter.convert(str(path))
        text = result.document.export_to_markdown()
        doc_stem = path.stem
        orig = text.strip()
        if not orig:
            return []

        parts = self._simple_split(orig, self._chunk_size, self._chunk_overlap)
        chunks = []
        for i, part in enumerate(parts):
            chunks.append(
                {
                    "chunk_id": f"{doc_stem}_{i}",
                    "source": "chunk",
                    "chunk_type": "text",
                    "chunk_text_original": part,
                    "chunk_text_embedded": f"[Document: {doc_stem}]\n{part}",
                    "page_number": None,
                    "section_title": None,
                    "language": _detect_language(part),
                    "chunk_hash": _chunk_hash(part),
                }
            )
        return chunks

    # ------------------------------------------------------------------
    # 5-stage pipeline
    # ------------------------------------------------------------------

    def _load_pipeline(self, path: Path) -> List[Dict]:
        doc_stem = path.stem

        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)

            # Stage 1 — parse
            elements, pic_info = self._stage1_parse(path, tmp_path)

            # Stage 2 — tables
            table_chunks = self._stage2_tables(elements, doc_stem)

            # Stage 3 — vision
            vision_map = self._stage3_vision(pic_info)

            # Stage 4 — text chunking (injects vision text at reading position)
            text_chunks = self._stage4_chunk(elements, vision_map, doc_stem)

            # Stage 5 — metadata enrichment + merge
            return self._stage5_metadata(text_chunks, table_chunks)

    # ------------------------------------------------------------------
    # Stage 1 — parse
    # ------------------------------------------------------------------

    def _stage1_parse(
        self, path: Path, tmp_dir: Path
    ) -> Tuple[List[Dict], List[Dict]]:
        """
        Returns
        -------
        elements:
            Flat list of dicts with keys: id, type, text, page, section_path,
            heading_level?, ordered?, png_path?
            Picture elements are included in-order with type="picture".
        pic_info:
            List of {id, png_path} for pictures that were saved as PNG files.
        """
        ext = path.suffix.lower()

        if ext == ".xlsx":
            elements = self._parse_excel(path)
            return elements, []

        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", message="Palette images with Transparency",
                                    category=UserWarning, module="PIL")
            result = self._converter.convert(str(path))
        doc = result.document

        elements, pic_info = self._extract_elements_and_pictures(doc, tmp_dir)

        # PPTX: enhance heading detection
        if ext == ".pptx":
            elements = self._reclassify_pptx_headings(path, elements)

        return elements, pic_info

    def _item_label(self, item: Any) -> Optional[str]:
        try:
            lbl = item.label
            if hasattr(lbl, "value"):
                return lbl.value
            return str(lbl)
        except AttributeError:
            return None

    def _page_no(self, item: Any) -> Optional[int]:
        try:
            return item.prov[0].page_no
        except (IndexError, AttributeError, TypeError):
            return None

    def _heading_level(self, item: Any) -> Optional[int]:
        try:
            lbl = self._item_label(item)
            if lbl in HEADING_LABELS:
                lvl = getattr(item, "level", None)
                if lvl is not None:
                    return int(lvl)
                return 1
        except Exception:
            pass
        return None

    def _is_ordered(self, item: Any) -> Optional[bool]:
        lbl = self._item_label(item)
        if lbl in LIST_LABELS:
            return bool(getattr(item, "enumerated", False) or getattr(item, "ordered", False))
        return None

    def _table_to_text(self, table_item: Any, doc: Any = None) -> str:
        try:
            md = table_item.export_to_markdown(doc=doc)
            if md:
                return md
        except Exception:
            pass
        try:
            return table_item.text or ""
        except Exception:
            return ""

    def _extract_elements_and_pictures(
        self, doc: Any, tmp_dir: Path
    ) -> Tuple[List[Dict], List[Dict]]:
        """Walk the document, building element list + extracting pictures to tmp_dir."""
        elements: List[Dict] = []
        pic_info: List[Dict] = []
        eid = 0
        section_stack: List[str] = []

        def _push_heading(text: str, level: int) -> None:
            nonlocal section_stack
            section_stack = [h for h in section_stack if _heading_depth(h) < level]
            section_stack.append(text)

        def _heading_depth(h: str) -> int:
            return section_stack.index(h) + 1 if h in section_stack else 1

        def _add(type_: str, text: str, page: Optional[int],
                 heading_level: Optional[int] = None,
                 ordered: Optional[bool] = None,
                 extra: Optional[Dict] = None) -> Dict:
            nonlocal eid
            el: Dict = {
                "id": f"e{eid}",
                "type": type_,
                "text": text,
                "page": page,
                "section_path": list(section_stack),
            }
            if heading_level is not None:
                el["heading_level"] = heading_level
            if ordered is not None:
                el["ordered"] = ordered
            if extra:
                el.update(extra)
            eid += 1
            return el

        try:
            for item, level in doc.iterate_items():
                text = ""
                try:
                    text = item.text or ""
                except AttributeError:
                    pass

                lbl = self._item_label(item)
                page = self._page_no(item)

                if lbl in HEADING_LABELS:
                    hl = self._heading_level(item) or level or 1
                    _push_heading(text, hl)
                    elements.append(_add("heading", text, page, heading_level=hl))

                elif lbl in PARAGRAPH_LABELS:
                    if text.strip():
                        elements.append(_add("paragraph", text, page))

                elif lbl in LIST_LABELS:
                    if text.strip():
                        marker = getattr(item, "marker", None)
                        if marker and self._is_ordered(item):
                            text = f"{marker.strip()} {text}"
                        elements.append(_add("list_item", text, page,
                                             ordered=self._is_ordered(item)))

                elif lbl in TABLE_LABELS:
                    ttext = self._table_to_text(item, doc=doc)
                    if ttext.strip():
                        elements.append(_add("table", ttext, page))

                elif lbl in PICTURE_LABELS:
                    # Track element ID before _add increments eid
                    pic_id = f"e{eid}"
                    png_path = self._save_picture(item, tmp_dir, pic_id)
                    pic_el = _add("picture", "", page, extra={"png_path": png_path})
                    elements.append(pic_el)
                    if png_path:
                        pic_info.append({"id": pic_id, "png_path": png_path})

        except AttributeError:
            logger.info("iterate_items not available, using doc.texts/tables")
            self._fallback_extract(doc, elements, section_stack, _add)

        # Fallback: also scan doc.pictures if iterate_items found none
        if not pic_info:
            pic_items = getattr(doc, "pictures", [])
            for idx, pic_item in enumerate(pic_items):
                page = self._page_no(pic_item)
                pic_id = f"epic{idx}"
                png_path = self._save_picture(pic_item, tmp_dir, pic_id)
                if png_path:
                    pic_el = {
                        "id": pic_id,
                        "type": "picture",
                        "text": "",
                        "page": page,
                        "section_path": [],
                        "png_path": png_path,
                    }
                    elements.append(pic_el)
                    pic_info.append({"id": pic_id, "png_path": png_path})

        return elements, pic_info

    def _fallback_extract(self, doc: Any, elements: List[Dict],
                          section_stack: List[str], _add) -> None:
        all_items = []
        for item in getattr(doc, "texts", []):
            all_items.append((item, self._page_no(item) or 0))
        for item in getattr(doc, "tables", []):
            all_items.append((item, self._page_no(item) or 0))
        all_items.sort(key=lambda x: x[1])
        for item, _ in all_items:
            lbl = self._item_label(item)
            page = self._page_no(item)
            try:
                text = item.text or ""
            except AttributeError:
                text = ""
            if lbl in HEADING_LABELS:
                hl = self._heading_level(item) or 1
                elements.append(_add("heading", text, page, heading_level=hl))
            elif lbl in PARAGRAPH_LABELS:
                if text.strip():
                    elements.append(_add("paragraph", text, page))
            elif lbl in LIST_LABELS:
                if text.strip():
                    marker = getattr(item, "marker", None)
                    if marker and self._is_ordered(item):
                        text = f"{marker.strip()} {text}"
                    elements.append(_add("list_item", text, page, ordered=self._is_ordered(item)))
            elif lbl in TABLE_LABELS:
                ttext = self._table_to_text(item, doc=doc)
                if ttext.strip():
                    elements.append(_add("table", ttext, page))

    def _save_picture(self, pic_item: Any, tmp_dir: Path, pic_id: str) -> Optional[str]:
        """Save a Docling picture item as PNG. Returns path or None."""
        try:
            img_obj = getattr(pic_item, "image", None)
            if img_obj is None:
                return None
            pil_image = getattr(img_obj, "pil_image", None)
            if pil_image is None and hasattr(img_obj, "as_pil"):
                pil_image = img_obj.as_pil()
            if pil_image is None:
                return None
            png_file = tmp_dir / f"{pic_id}.png"
            pil_image.save(str(png_file))
            return str(png_file)
        except Exception as exc:
            logger.debug("Could not save picture %s: %s", pic_id, exc)
            return None

    def _parse_excel(self, path: Path) -> List[Dict]:
        """Parse XLSX with openpyxl → flat element list."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required for Excel support: pip install openpyxl")

        wb = openpyxl.load_workbook(str(path), data_only=True)
        elements: List[Dict] = []
        eid = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            elements.append({
                "id": f"e{eid}", "type": "heading", "text": sheet_name,
                "page": None, "section_path": [], "heading_level": 1,
            })
            eid += 1
            rows: List[List[str]] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(cells)
            if not rows:
                continue
            header = rows[0]
            sep = ["---"] * len(header)
            md_lines = [
                "| " + " | ".join(header) + " |",
                "| " + " | ".join(sep) + " |",
            ]
            for row in rows[1:]:
                padded = row + [""] * max(0, len(header) - len(row))
                md_lines.append("| " + " | ".join(padded[:len(header)]) + " |")
            elements.append({
                "id": f"e{eid}", "type": "table",
                "text": "\n".join(md_lines),
                "page": None, "section_path": [sheet_name],
            })
            eid += 1
        return elements

    def _reclassify_pptx_headings(self, path: Path, elements: List[Dict]) -> List[Dict]:
        """Use python-pptx signals to promote eligible shapes to headings."""
        try:
            heading_texts = self._pptx_heading_texts(str(path))
        except Exception as exc:
            logger.debug("PPTX heading reclassification skipped: %s", exc)
            return elements
        if not heading_texts:
            return elements

        def _norm(text: str) -> str:
            return re.sub(r'\s+', ' ', text.strip()).lower()

        heading_set: set = set()
        for page_set in heading_texts.values():
            for t in page_set:
                heading_set.add(_norm(t))

        result = []
        for el in elements:
            if el.get("type") == "paragraph" and _norm(el.get("text", "")) in heading_set:
                el = dict(el)
                el["type"] = "heading"
                el["heading_level"] = el.get("heading_level", 1)
            result.append(el)
        return result

    def _pptx_heading_texts(self, pptx_path: str) -> Dict[int, set]:
        """Detect heading shapes in PPTX via bold/font-size/color signals."""
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return {}

        def _collect_runs(shape) -> list:
            runs = []
            if shape.shape_type == 6:
                try:
                    for child in shape.shapes:
                        runs.extend(_collect_runs(child))
                except Exception:
                    pass
                return runs
            if not shape.has_text_frame:
                return runs
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        runs.append(run)
            return runs

        def _text_of(shape) -> str:
            try:
                if shape.shape_type == 6:
                    lines = []
                    for child in shape.shapes:
                        lines.append(_text_of(child))
                    return " ".join(lines)
                return shape.text_frame.text if shape.has_text_frame else ""
            except Exception:
                return ""

        def _is_heading(shape) -> bool:
            runs = _collect_runs(shape)
            if not runs:
                return False
            if all(getattr(r.font, "bold", False) for r in runs):
                return True
            sizes = [r.font.size for r in runs if r.font.size is not None]
            if sizes:
                median_size = sorted(sizes)[len(sizes) // 2]
                if all(s >= median_size * 1.25 for s in sizes if s):
                    return True
            return False

        prs = Presentation(pptx_path)
        result: Dict[int, set] = {}
        for slide_num, slide in enumerate(prs.slides, start=1):
            headings: set = set()
            for shape in slide.shapes:
                if _is_heading(shape):
                    t = _text_of(shape).strip()
                    if t:
                        headings.add(t)
            if headings:
                result[slide_num] = headings
        return result

    # ------------------------------------------------------------------
    # Stage 2 — table processing
    # ------------------------------------------------------------------

    def _stage2_tables(self, elements: List[Dict], doc_stem: str) -> List[Dict]:
        """Extract and process all table elements → table chunk dicts."""
        table_chunks: List[Dict] = []
        tidx = 0

        for el in elements:
            text = el.get("text", "")
            is_table = el.get("type") == "table"
            is_ascii_table = (
                el.get("type") in ("paragraph", "text") and
                _has_heavy_box_drawing(text)
            )
            if not (is_table or is_ascii_table) or not text.strip():
                continue

            section_path = el.get("section_path", [])
            section_title = section_path[-1] if section_path else None
            page = el.get("page")
            rows, cols = _estimate_table_dimensions(text)
            diff = _ocr_difficulty(rows, cols)
            table_type = _classify_table_type(text, rows, cols)

            if table_type == "faq":
                headers, data_rows = _parse_faq_rows(text)
                ridx = 0
                for row in data_rows:
                    row_text = _faq_row_to_text(headers, row)
                    if not row_text.strip():
                        continue
                    table_chunks.append({
                        "chunk_id": f"{doc_stem}_table_{tidx}_row_{ridx}",
                        "text": row_text,
                        "text_natural": None,
                        "page_number": page,
                        "section_title": section_title,
                        "ocr_difficulty": diff,
                        "rows": rows,
                        "cols": cols,
                        "table_type": table_type,
                    })
                    ridx += 1
                logger.info("Table %d (faq): %d row chunks", tidx, ridx)
            else:
                text_natural = self._narrate_table(text, section_title)
                table_chunks.append({
                    "chunk_id": f"{doc_stem}_table_{tidx}",
                    "text": text,
                    "text_natural": text_natural,
                    "page_number": page,
                    "section_title": section_title,
                    "ocr_difficulty": diff,
                    "rows": rows,
                    "cols": cols,
                    "table_type": table_type,
                })
                logger.info("Table %d (%s): narrated=%s", tidx, table_type, bool(text_natural))
            tidx += 1

        return table_chunks

    def _narrate_table(self, markdown_text: str, section_title: Optional[str]) -> Optional[str]:
        """Call Ollama to convert a table to natural language. Returns None on failure."""
        prompt = _NARRATE_PROMPT_TMPL.format(
            section_title=section_title or "(unknown)",
            markdown_text=markdown_text,
        )
        result = self._ollama_text(prompt)
        if result:
            result = _strip_preamble(result)
            logger.debug("Table narration: %d chars", len(result))
        return result or None

    # ------------------------------------------------------------------
    # Stage 3 — vision summarization
    # ------------------------------------------------------------------

    def _stage3_vision(self, pic_info: List[Dict]) -> Dict[str, str]:
        """Summarize each extracted picture via 3-pass Ollama VLM."""
        vision_map: Dict[str, str] = {}
        for pi in pic_info:
            png_path = pi.get("png_path")
            if not png_path or not Path(png_path).exists():
                continue
            summary = self._summarize_image(png_path)
            if summary:
                vision_map[pi["id"]] = summary
                logger.info("Vision: %s → %d chars", pi["id"], len(summary))
        return vision_map

    def _summarize_image(self, image_path: str) -> str:
        """3-pass Ollama pipeline: classify → structured extraction → synthesis."""
        try:
            b64 = self._prepare_image(image_path)
        except Exception as exc:
            logger.warning("Could not load image %s: %s", image_path, exc)
            return ""

        # Pass 1: image type classification
        caption = self._ollama_vision(self._vlm_model, _CAPTION_PROMPT, b64, max_tokens=60)
        if not caption:
            return ""
        type_key = _classify_vision_caption(caption)
        type_key = _size_guard(image_path, type_key)
        logger.debug("Vision pass1: type=%s caption=%r", type_key, caption[:80])

        # Pass 2: structured extraction
        pass2_prompt = _STRUCTURED_PROMPTS.get(type_key, _DEFAULT_STRUCTURED_PROMPT)
        ocr = self._ollama_vision(self._vlm_model, pass2_prompt, b64, max_tokens=2000)
        if _is_ocr_echo(ocr):
            logger.warning("Vision pass2: OCR echo detected, discarding")
            ocr = ""

        # Pass 3: text-only synthesis
        if caption.strip() or ocr.strip():
            ocr_deduped = _dedup_ocr(ocr)
            synthesis_prompt = _SYNTHESIS_PROMPT.format(
                caption=caption.strip() or "(not available)",
                ocr_deduped=ocr_deduped.strip() or "(not available)",
            )
            synthesized = self._ollama_text(synthesis_prompt)
            if synthesized and synthesized.strip():
                return synthesized.strip()

        # Fallback: caption + OCR concatenated
        out = []
        if caption:
            out.append(f"[Caption] {caption}")
        if ocr:
            out.append(f"[OCR] {ocr}")
        return "\n".join(out)

    def _prepare_image(self, image_path: str) -> str:
        """Return base64-encoded PNG, upscaling small images for better VLM results."""
        _MIN_W, _MIN_H = 400, 150
        from PIL import Image as _PILImage
        with _PILImage.open(image_path) as img:
            w, h = img.size
            scale = max(
                _MIN_W / w if w < _MIN_W else 1.0,
                _MIN_H / h if h < _MIN_H else 1.0,
            )
            if scale > 1.0:
                new_w, new_h = int(w * scale), int(h * scale)
                img = img.resize((new_w, new_h), _PILImage.LANCZOS)
                logger.info("Upscaled %dx%d → %dx%d for VLM", w, h, new_w, new_h)
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="PNG")
            return base64.b64encode(buf.getvalue()).decode()

    # ------------------------------------------------------------------
    # Stage 4 — text chunking
    # ------------------------------------------------------------------

    def _stage4_chunk(
        self,
        elements: List[Dict],
        vision_map: Dict[str, str],
        doc_stem: str,
    ) -> List[Dict]:
        """Split text elements into chunks with contextual prefix."""
        rcts = _build_rcts(self._chunk_size, self._chunk_overlap)
        groups = _group_elements(elements, vision_map)
        raw_chunks: List[Dict] = []
        chunk_idx = 0

        for group in groups:
            if not group:
                continue
            if len(group) == 1 and group[0].get("type") == "heading":
                continue  # heading-only chunks add no retrievable content

            first = group[0]
            page = first.get("page")
            section_path = first.get("section_path", [])
            section_title = section_path[-1] if section_path else None
            block_text = "\n".join(
                el.get("text", "") for el in group if el.get("text", "").strip()
            )
            source_ids = [el["id"] for el in group]
            ptype = _preliminary_type(first)

            if not block_text.strip():
                continue

            if rcts and len(block_text) > self._chunk_size and ptype == "text":
                sub_texts = rcts.split_text(block_text)
            else:
                sub_texts = self._simple_split(block_text, self._chunk_size, self._chunk_overlap)

            for sub_text in sub_texts:
                if not sub_text.strip():
                    continue
                raw_chunks.append({
                    "chunk_id": f"{doc_stem}_{chunk_idx}",
                    "source_element_ids": source_ids,
                    "text": sub_text,
                    "page_number": page,
                    "section_title": section_title,
                    "preliminary_type": ptype,
                })
                chunk_idx += 1

        # Post-processing
        chunks = _merge_short_text_chunks(raw_chunks)
        chunks = _merge_warning_headers(chunks)

        # Add contextual prefix (Contextual Retrieval)
        for chunk in chunks:
            orig = chunk.pop("text")
            prefix = _build_context_prefix(doc_stem, chunk.get("section_title"), chunk.get("page_number"))
            chunk["chunk_text_original"] = orig
            chunk["chunk_text_embedded"] = f"{prefix}\n{orig}"

        return chunks

    @staticmethod
    def _simple_split(text: str, size: int, overlap: int) -> List[str]:
        chunks: List[str] = []
        start = 0
        while start < len(text):
            end = min(start + size, len(text))
            chunks.append(text[start:end])
            if end == len(text):
                break
            start += size - overlap
        return chunks

    # ------------------------------------------------------------------
    # Stage 5 — metadata enrichment + merge
    # ------------------------------------------------------------------

    def _stage5_metadata(
        self,
        text_chunks: List[Dict],
        table_chunks: List[Dict],
    ) -> List[Dict]:
        """Merge text + table chunks, assign chunk_type, language, hash."""
        unified: List[Dict] = []

        # --- text / procedure / picture chunks from stage 4 ---
        for c in text_chunks:
            orig = c.get("chunk_text_original", "")
            embedded = c.get("chunk_text_embedded", orig)
            ptype = c.get("preliminary_type", "text")
            if ptype == "picture":
                chunk_type = "picture"
            elif ptype == "procedure_step":
                chunk_type = "procedure_step"
            else:
                chunk_type = "text"
            unified.append({
                "chunk_id": c["chunk_id"],
                "source": "chunk",
                "chunk_type": chunk_type,
                "chunk_text_original": orig,
                "chunk_text_embedded": embedded,
                "page_number": c.get("page_number"),
                "section_title": c.get("section_title"),
                "language": _detect_language(orig),
                "chunk_hash": _chunk_hash(orig),
            })

        # --- table chunks from stage 2 ---
        for tc in table_chunks:
            orig = tc.get("text_natural") or tc.get("text", "")
            markdown = tc.get("text", "")
            # Derive doc_stem from chunk_id (e.g. "mystem_table_0" → "mystem")
            doc_stem_guess = tc["chunk_id"].split("_table_")[0]
            prefix = _build_context_prefix(
                doc_stem_guess,
                tc.get("section_title"),
                tc.get("page_number"),
            )
            embedded = f"{prefix}\n{orig}"
            entry: Dict = {
                "chunk_id": tc["chunk_id"],
                "source": "table",
                "chunk_type": "parts_table",
                "chunk_text_original": orig,
                "chunk_text_embedded": embedded,
                "chunk_text_raw": markdown,
                "page_number": tc.get("page_number"),
                "section_title": tc.get("section_title"),
                "language": _detect_language(orig),
                "chunk_hash": _chunk_hash(orig),
                "ocr_difficulty": tc.get("ocr_difficulty"),
                "rows": tc.get("rows"),
                "cols": tc.get("cols"),
                "table_type": tc.get("table_type"),
            }
            unified.append(entry)

        logger.info(
            "Stage 5: %d text/proc + %d table = %d total chunks",
            len(text_chunks), len(table_chunks), len(unified),
        )
        return unified

    # ------------------------------------------------------------------
    # Ollama HTTP helpers
    # ------------------------------------------------------------------

    def _ollama_text(self, prompt: str, max_tokens: int = 600) -> str:
        """Text-only Ollama /api/generate call (uses INDEX_MODEL / text_model)."""
        return self._ollama_call(
            model=self._text_model,
            prompt=prompt,
            image_b64=None,
            max_tokens=max_tokens,
        )

    def _ollama_vision(self, model: str, prompt: str, image_b64: str, max_tokens: int = 800) -> str:
        """Vision Ollama /api/generate call with image attachment."""
        return self._ollama_call(
            model=model,
            prompt=prompt,
            image_b64=image_b64,
            max_tokens=max_tokens,
        )

    def _ollama_call(
        self, model: str, prompt: str, image_b64: Optional[str], max_tokens: int
    ) -> str:
        import urllib.error
        import urllib.request

        payload: Dict = {
            "model": model,
            "prompt": prompt,
            "stream": False,
            "options": {"temperature": 0.1, "num_predict": max_tokens},
        }
        if image_b64:
            payload["images"] = [image_b64]

        data = json.dumps(payload).encode()
        req = urllib.request.Request(
            f"{self._base_url}/api/generate",
            data=data,
            headers={"Content-Type": "application/json"},
        )
        try:
            with urllib.request.urlopen(req, timeout=self._timeout) as resp:
                result = json.loads(resp.read())
                return result.get("response", "").strip()
        except Exception as exc:
            logger.warning("Ollama call failed (%s): %s", self._base_url, exc)
            return ""

    # ------------------------------------------------------------------
    # Converter builder (picture description disabled — handled in Stage 3)
    # ------------------------------------------------------------------
    # Converter builder (picture description disabled — handled in Stage 3)
    # ------------------------------------------------------------------

    def _build_converter(self):
        """Build Docling DocumentConverter with picture description disabled."""
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import PdfPipelineOptions
        from docling.document_converter import (
            DocumentConverter,
            ImageFormatOption,
            PdfFormatOption,
            WordFormatOption,
        )

        pdf_opts = PdfPipelineOptions()
        pdf_opts.do_ocr = False
        pdf_opts.generate_picture_images = True

        return DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pdf_opts),
                InputFormat.IMAGE: ImageFormatOption(pipeline_options=pdf_opts),
                InputFormat.DOCX: WordFormatOption(pipeline_options=pdf_opts),
            }
        )


